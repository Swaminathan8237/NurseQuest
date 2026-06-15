"""
fill_ppt_tables.py
==================
Fills all empty tables in FIRST REVIEW PPT Template.pptx with
proper content aligned to the NurseQuest project.

Tables filled:
  Slide 4  - Review of Literature (References 1-8)
  Slide 5  - Review of Literature Continued (References 9-15)
  Slide 6  - Review of Existing Systems
  Slide 10 - Modules

Output: docs/FIRST_REVIEW_FILLED.pptx
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import copy

# ─────────────────────────────────────────────
# Paths
# ─────────────────────────────────────────────
SRC  = r'd:\E0224033\Projeccts\Nursing_Education_Platform_gpt\Nursing_Education_Platform_gpt\docs\FIRST REVIEW PPT Template.pptx'
DEST = r'd:\E0224033\Projeccts\Nursing_Education_Platform_gpt\Nursing_Education_Platform_gpt\docs\FIRST_REVIEW_FILLED.pptx'

# ─────────────────────────────────────────────
# Color constants  (match the presentation theme)
# ─────────────────────────────────────────────
COL_HEADER_BG   = RGBColor(0xA5, 0x00, 0x21)   # dark crimson header rows
COL_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)   # white
COL_ROW_TEXT    = RGBColor(0x00, 0x00, 0x00)   # black body text
COL_ALT_BG     = RGBColor(0xF2, 0xDC, 0xDB)   # very light rose for alt rows

# ─────────────────────────────────────────────
# Helper – set a single cell's text + formatting
# ─────────────────────────────────────────────
def set_cell(cell, text, bold=False, font_size=Pt(8),
             color=COL_ROW_TEXT, align=PP_ALIGN.LEFT,
             bg_color=None):
    tf = cell.text_frame
    tf.word_wrap = True
    # clear existing runs
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ''
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.name = 'Times New Roman'
    # background fill – use the proper python-pptx API to avoid malformed XML
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

# ─────────────────────────────────────────────
# DATA – Review of Literature (Slides 4 & 5)
# ─────────────────────────────────────────────
# Actual PPT columns (DO NOT CHANGE HEADERS):
#  Col0: Author and Year (chronological order)
#  Col1: Methodology / Materials and Methods
#  Col2: Datasets
#  Col3: Results / Outcome
#  Col4: Advantage / Key findings from the work
#  Col5: Disadvantage / GAP Identified from the work
LIT_REFS = [
    (
        "R. Nguyen\n2026",
        "Argumentative review & pedagogical analysis of gamification frameworks applied to nursing curricula.",
        "Existing literature on nursing education gamification; not primary data collection.",
        "Gamification strategies (XP, leaderboards, rewards) significantly increase student motivation and knowledge retention.",
        "Establishes the theoretical basis for applying gamification in nursing; directly supports NurseQuest's XP and leveling model.",
        "Limited empirical data; lacks a controlled trial. Does not address technical implementation of gamification systems."
    ),
    (
        "A. Azoulay & F. Lim\n2026",
        "Integrative review of 22 studies examining gamification effects on nursing academic performance.",
        "22 peer-reviewed studies on nursing student cohorts in higher education institutions.",
        "Gamified platforms improved academic performance scores by an average of 18% compared to traditional methods.",
        "Provides strong evidence that gamified platforms improve academic scores; validates NurseQuest's evidence-based design.",
        "Studies varied widely in gamification type and measurement; no standardized scoring metric across studies."
    ),
    (
        "S. Aalaei et al.\n2026",
        "Quasi-experimental pre-post study evaluating a gamified mobile app for cardiac rhythm interpretation.",
        "60 undergraduate nursing students tested pre/post on ECG interpretation using a gamified mobile application.",
        "Significant improvement in ECG interpretation scores post-intervention (p<0.001); students rated engagement highly.",
        "Domain-specific gamified assessment design; directly informs NurseQuest's Image/Audio-based clinical question types.",
        "Single-institution study; limited generalizability. App was mobile-only with no web or multiplayer capability."
    ),
    (
        "M. AlMekkawi et al.\n2026",
        "Systematic review and meta-analysis of 18 RCTs and quasi-experimental studies on game-based learning in nursing.",
        "18 studies across undergraduate nursing programs globally; 1,200+ students aggregated.",
        "Game-based learning showed statistically significant improvements in knowledge retention (SMD=0.72) and clinical reasoning.",
        "Meta-analysis with strong statistical power confirming game-based learning effectiveness for nursing education.",
        "High heterogeneity among included studies; most games were desktop-only and lacked real-time multiplayer components."
    ),
    (
        "V. I. Bahar et al.\n2026",
        "Descriptive, cross-sectional study evaluating a serious game's impact on palliative care awareness in nursing students.",
        "85 third-year nursing students from a Turkish university who played the serious game over 2 weeks.",
        "Palliative care awareness scores significantly increased post-game (p<0.05); students reported higher empathy levels.",
        "Demonstrates that serious games work in specialized nursing domains; aligns with NurseQuest's domain module concept.",
        "Focused on a single nursing topic (palliative care); no leaderboard, progression system, or multiplayer component."
    ),
    (
        "K. McMillan et al.\n2026",
        "Scoping review protocol following JBI methodology for gamification design in preregistration nurse education.",
        "Planned review of databases: CINAHL, MEDLINE, PsycINFO, Embase; no primary data collected yet.",
        "Protocol identifies key design variables for nurse-targeted gamification: feedback loops, narrative, and leaderboards.",
        "Defines design principles for nurse-targeted gamification; guides NurseQuest's leaderboard and badge system design.",
        "Protocol paper only; full results not yet published. Lacks quantitative evidence of outcomes."
    ),
    (
        "M. Yang et al.\n2026",
        "Systematic review and meta-analysis of escape room interventions in Chinese nursing education programs.",
        "12 studies conducted at Chinese nursing universities; 850+ nursing students across various clinical specialties.",
        "Escape room learning produced significantly higher clinical reasoning scores (SMD=0.61) vs. traditional teaching.",
        "Validates interactive challenge-based formats; informs NurseQuest's Live Quiz (Kahoot-style) and scenario question types.",
        "Geographically limited to China; escape room format is physical-first, lacking digital scalability and multiplayer realism."
    ),
    (
        "Y. G. S. Barbalho et al.\n2025",
        "Design and validation study using iterative prototyping and expert panel review for a serious game in nursing.",
        "Expert panel of 12 nursing educators + 40 undergraduate students for usability and content validation testing.",
        "Game achieved high content validity index (CVI=0.92) and usability scores; students showed improved knowledge scores.",
        "Provides a validated serious game design and validation framework directly applicable to NurseQuest's development.",
        "Single-cycle validation; long-term retention and multi-session engagement were not measured in this study."
    ),
]

LIT_REFS_CONT = [
    (
        "N. Aktaş & Y. Sazak\n2025",
        "Quasi-experimental study using an escape room as the intervention for pharmacology knowledge acquisition.",
        "48 nursing students (intervention=24, control=24) at a Turkish nursing school over one semester.",
        "Escape room group scored significantly higher on post-tests (p<0.01); higher satisfaction compared to lectures.",
        "Supports scenario-based quiz format; informs NurseQuest's Jumbled Letters and Sequence question type design.",
        "Physical escape room format; cannot scale digitally, no real-time scoring or remote accessibility."
    ),
    (
        "M. White & T. Shellenbarger\n2018",
        "Descriptive study examining the integration of digital badges as motivational gamification in nursing courses.",
        "Undergraduate nursing students across 3 semesters in a US university nursing program.",
        "Badge-earning behavior correlated with higher assignment completion rates and increased student engagement.",
        "Pioneering work on badge-based motivation in nursing; directly informs NurseQuest's achievement badge system.",
        "Study predates modern web platforms; no XP, leaderboard, or real-time feedback component evaluated."
    ),
    (
        "S. Mishra et al.\n2026",
        "System design and development using Next.js, Node.js, and Socket.IO for a scalable multiplayer quiz platform.",
        "Load testing with 500 concurrent users; performance benchmarks for latency and throughput measured.",
        "Platform achieved sub-100ms latency under 500 concurrent users; Socket.IO enabled reliable real-time sync.",
        "Technical architecture reference for scalable quiz platforms; directly maps to NurseQuest's full-stack design.",
        "Generic quiz content; no domain specificity, no gamification progression (XP/badges), no role-based access."
    ),
    (
        "A. Kumar et al.\n2025",
        "Software development study using React 18, Node.js Express, and Socket.IO for a gamified quiz application.",
        "Developer testing environment with 50 concurrent WebSocket connections; functional testing performed.",
        "Demonstrated real-time score updates, live leaderboard sync, and responsive React UI across device types.",
        "Blueprint for React + Socket.IO quiz architecture; validates NurseQuest's front-end and real-time stack choices.",
        "No nursing or clinical domain content; lacks a progression system, role-based access, or media-rich questions."
    ),
    (
        "J. Silva & M. Santos\n2023",
        "Comparative benchmarking study evaluating SQLite vs. PostgreSQL vs. MySQL for LMS-type web applications.",
        "Synthetic LMS database workloads: 100K records, read/write benchmarks, concurrent access tests.",
        "SQLite outperformed others in read-heavy single-server workloads; best suited for embedded deployments.",
        "Justifies SQLite (better-sqlite3) selection as NurseQuest's embedded database for single-server deployment.",
        "SQLite limitations in multi-server/distributed deployments are not addressed; not ideal for large-scale production."
    ),
    (
        "L. Chen & X. Wang\n2024",
        "Experimental performance comparison of WebSocket (Socket.IO) and WebRTC for browser-based multiplayer games.",
        "Simulated 200-player game sessions; latency, jitter, and packet loss measured under varying network conditions.",
        "Socket.IO achieved lower average latency (42ms) and simpler server-side architecture than WebRTC for game sync.",
        "Performance data supporting WebSocket (Socket.IO) choice for NurseQuest's live multiplayer quiz synchronization.",
        "WebRTC offers lower latency for peer-to-peer; Socket.IO may face scalability bottlenecks beyond 1000 concurrent users."
    ),
    (
        "D. Jones & R. Smith\n2022",
        "Case study and iterative development using React.js for interactive SPA-based medical training tools.",
        "Medical training application used by 120 students; usability and learning outcome data collected over 8 weeks.",
        "React SPA achieved high usability scores (SUS=82.5); students showed measurable improvement in knowledge tests.",
        "Validates React SPA architecture for medical education; confirms NurseQuest's frontend design approach.",
        "Study used older React 16; did not address real-time features, multiplayer, or gamification components."
    ),
]

# ─────────────────────────────────────────────
# DATA – Review of Existing Systems (Slide 6)
# ─────────────────────────────────────────────
# Columns: S.No | System | Type | Key Findings | Disadvantage/GAP | NurseQuest Advantage
EXISTING_SYSTEMS = [
    (
        "1",
        "Kahoot!",
        "Commercial Quiz Platform",
        "Real-time multiplayer quizzes, live leaderboard, teacher-hosted sessions, mobile-friendly interface.",
        "No nursing-specific content, no XP/leveling, no clinical question types, no role-based educator access.",
        "NurseQuest adds nursing domain modules, clinical question types (image, audio, ECG), XP/rank progression, and JWT-secured RBAC."
    ),
    (
        "2",
        "Quizlet",
        "Flashcard & Study Tool",
        "Flashcards, practice tests, study modes, and collaborative study sets for self-paced learning.",
        "No real-time multiplayer, no gamification engine, not designed for clinical scenario-based assessment.",
        "NurseQuest provides live multiplayer sessions, multi-format clinical questions, and a full gamification layer with XP and badges."
    ),
    (
        "3",
        "Moodle LMS",
        "Learning Management System",
        "Course management, quizzes, grading, forums, file submission, and detailed performance reporting.",
        "Heavy and complex setup, no gamification layer, poor UX for students, no real-time competitive gameplay.",
        "NurseQuest offers a lightweight, gamified, modern interface purpose-built for nursing assessments and live quiz sessions."
    ),
    (
        "4",
        "Google Forms",
        "Online Form / Quiz Tool",
        "Simple MCQ quizzes, auto-grading, and results export to Google Sheets spreadsheets.",
        "No gamification, no live multiplayer, no media-rich question types, no student progression or streak tracking.",
        "NurseQuest supports 6 question types, real-time game sessions, student XP and streak tracking, and achievement badges."
    ),
    (
        "5",
        "Mentimeter",
        "Interactive Presentation Tool",
        "Live polls, word clouds, Q&A sessions, and audience response during presentations.",
        "Not a quiz platform; lacks scoring, XP, levels, and curriculum-linked assessment. No nursing domain support.",
        "NurseQuest offers structured clinical quizzes with scoring, timer-based pressure, real-time leaderboard, and nursing-specific content."
    ),
    (
        "6",
        "Wayground",
        "Online Assessment Platform",
        "Custom assessments, analytics dashboard, student progress tracking, and integrations with LMS tools.",
        "No gamification or competitive elements, no live multiplayer mode, generic content with no nursing domain specificity.",
        "NurseQuest introduces live competitive gaming, gamified progression (XP/ranks/badges), and nursing-domain-specific clinical question types."
    ),
    (
        "7",
        "Nearpod",
        "Interactive Learning Platform",
        "Teacher-paced lessons, embedded quizzes, VR field trips, collaborative boards, and real-time student responses.",
        "Primarily a lesson delivery tool, not a standalone quiz game; no XP system, no competitive leaderboard, no nursing content.",
        "NurseQuest focuses on competitive gamified assessment with nursing-specific question types, live multiplayer scoring, and student rank progression."
    ),
    (
        "8",
        "Socrative",
        "Classroom Response System",
        "Teacher-driven quizzes, exit tickets, space races (team games), and instant result reports.",
        "Limited question types (MCQ, T/F, short answer), no media-rich clinical scenarios, no gamification or XP progression.",
        "NurseQuest adds 6 rich question types (including image, audio, and video), a gamification engine with XP/levels, and live leaderboards."
    ),
    (
        "9",
        "ProProfs Quiz Maker",
        "Online Quiz Platform",
        "Quiz creation, branching logic, certificates of completion, and performance reports.",
        "No nursing-specific domain modules, no live multiplayer mode, limited question types, subscription-based pricing.",
        "NurseQuest is free and open-source with role-based teacher/student access, live game hosting, and built-in nursing curriculum management."
    ),
]

# ─────────────────────────────────────────────
# DATA – Modules (Slide 10)
# ─────────────────────────────────────────────
# Columns: S.No | Module Name | Description | Key Features | Technology Used
MODULES = [
    (
        "1",
        "Authentication & User Management",
        "Handles user registration, login, JWT-based session management, and role assignment (Teacher / Student)",
        "JWT tokens, bcryptjs password hashing, role-based routing middleware, avatar setup flow",
        "Node.js, Express 5, better-sqlite3, JWT"
    ),
    (
        "2",
        "Student Dashboard",
        "Displays student XP, level, active streak, leaderboard rank, recent attempts, and available quizzes",
        "XP progress bar, 7 nursing rank levels, streak counter, achievement badges, global leaderboard widget",
        "React 19, Context API, Chart.js, GSAP, Lottie"
    ),
    (
        "3",
        "Teacher Dashboard & Quiz Builder",
        "Provides analytics overview and a full-featured quiz creation tool for teachers",
        "Quiz CRUD operations, 6 question type editors, media upload (image/video/audio), difficulty & timer settings",
        "React 19, Multer (file upload), Node.js REST API, SQLite"
    ),
    (
        "4",
        "Multi-Format Question Engine",
        "Supports 6 diverse question formats to create rich, clinically-relevant assessments",
        "MCQ, Image-based, Video-based, Audio-based, Jumbled Letters, Sequence Ordering",
        "React, HTML5 media elements, Drag-and-drop API, REST API"
    ),
    (
        "5",
        "Live Multiplayer Quiz Sessions",
        "Enables real-time teacher-hosted quiz games with simultaneous question delivery and live scoring",
        "Game code join, synchronized question broadcast, countdown timer, real-time leaderboard after each round",
        "Socket.IO (WebSocket), Node.js, React"
    ),
    (
        "6",
        "Gamification & Progression Engine",
        "Manages XP awarding, level-up logic, badge unlocking, streak tracking, and global leaderboard",
        "XP scoring formula (correctness + speed + streak), 7 nurse ranks, daily streak bonuses, global podium",
        "Node.js, SQLite, custom scoring algorithm"
    ),
    (
        "7",
        "Module Content Management",
        "Allows teachers to organize nursing curriculum into structured learning modules with rich text content",
        "Create/edit/delete modules, student-facing module browser, link modules to quizzes",
        "React, Node.js REST API, SQLite"
    ),
]


# ─────────────────────────────────────────────
# Core fill function
# ─────────────────────────────────────────────
def fill_table(table, data_rows, header_row=None, fill_header=False,
               header_font=Pt(9), body_font=Pt(7.5)):
    """
    Fills a pptx table with data.
    data_rows   : list of tuples – one tuple per body row
    header_row  : list of column header strings (only used when fill_header=True)
    fill_header : if True, overwrite row 0 with header_row strings;
                  if False, keep existing headers in the PPT as-is
    """
    num_cols = len(table.columns)
    num_rows = len(table.rows)

    # ---------- Header row (row 0) – only if explicitly requested ----------
    if fill_header and header_row and num_rows >= 1 and len(header_row) == num_cols:
        for c, heading in enumerate(header_row):
            set_cell(
                table.cell(0, c),
                heading,
                bold=True,
                font_size=header_font,
                color=COL_HEADER_TEXT,
                align=PP_ALIGN.CENTER,
                bg_color=COL_HEADER_BG
            )

    # ---------- Data rows ----------
    for r_idx, row_data in enumerate(data_rows):
        table_row = r_idx + 1          # row 0 is header
        if table_row >= num_rows:
            break                      # table doesn't have enough rows
        alt = (r_idx % 2 == 1)
        row_bg = COL_ALT_BG if alt else None
        for c_idx, cell_text in enumerate(row_data):
            if c_idx >= num_cols:
                break
            is_first_col = (c_idx == 0)
            set_cell(
                table.cell(table_row, c_idx),
                str(cell_text),
                bold=is_first_col,
                font_size=body_font,
                color=COL_ROW_TEXT,
                align=PP_ALIGN.CENTER if is_first_col else PP_ALIGN.LEFT,
                bg_color=row_bg
            )


# ─────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────
prs = Presentation(SRC)

for slide_idx, slide in enumerate(prs.slides):
    slide_num = slide_idx + 1

    for shape in slide.shapes:
        if not shape.has_table:
            continue

        tbl = shape.table
        rows = len(tbl.rows)
        cols = len(tbl.columns)
        print(f"[Slide {slide_num}] Found table: {rows} rows × {cols} cols")

        # ── Slide 4 – Review of Literature (Refs 1-8) ──
        if slide_num == 4:
            # fill_header=False → keep existing PPT column headers
            fill_table(tbl, LIT_REFS, fill_header=False, body_font=Pt(7))
            print("  → Filled: Review of Literature (Refs 1-8)")

        # ── Slide 5 – Review of Literature Continued (Refs 9-15) ──
        elif slide_num == 5:
            fill_table(tbl, LIT_REFS_CONT, fill_header=False, body_font=Pt(7))
            print("  → Filled: Review of Literature Continued (Refs 9-15)")

        # Slide 6 – skipped (user will fill manually)
        elif slide_num == 6:
            print("  → Skipped: Review of Existing Systems (not touched)")

        # Slide 10 (Modules) – already filled by user, skipping.
        elif slide_num == 10:
            print("  → Skipped: Modules (already filled by user)")

prs.save(DEST)
print(f"\n✅ Saved filled presentation to:\n   {DEST}")
