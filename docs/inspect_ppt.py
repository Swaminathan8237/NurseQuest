"""Inspect the PPT template structure and then fill it with problem statement content."""
import subprocess
import sys

# Ensure python-pptx is installed
subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import os

ppt_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "FIRST REVIEW PPT Template.pptx")

prs = Presentation(ppt_path)

print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")
print(f"Number of slides: {len(prs.slides)}")
print(f"Number of slide layouts: {len(prs.slide_layouts)}")
print()

for i, slide in enumerate(prs.slides):
    print(f"=== Slide {i+1} ===")
    print(f"  Layout: {slide.slide_layout.name}")
    for j, shape in enumerate(slide.shapes):
        print(f"  Shape {j}: name='{shape.shape_id}:{shape.name}', type={shape.shape_type}")
        print(f"    Position: left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}")
        if shape.has_text_frame:
            for k, para in enumerate(shape.text_frame.paragraphs):
                text = para.text
                if text.strip():
                    print(f"    Paragraph {k}: '{text[:80]}...' " if len(text)>80 else f"    Paragraph {k}: '{text}'")
                    if para.runs:
                        run = para.runs[0]
                        print(f"      Font: name={run.font.name}, size={run.font.size}, bold={run.font.bold}, color={run.font.color.rgb if run.font.color and run.font.color.rgb else 'None'}")
                else:
                    print(f"    Paragraph {k}: (empty)")
    print()

print("\n=== Available Slide Layouts ===")
for i, layout in enumerate(prs.slide_layouts):
    print(f"  Layout {i}: '{layout.name}'")
    for ph in layout.placeholders:
        print(f"    Placeholder idx={ph.placeholder_format.idx}, name='{ph.name}', type={ph.placeholder_format.type}")
