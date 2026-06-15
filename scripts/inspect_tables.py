from pptx import Presentation

pptx_path = r'd:\E0224033\Projeccts\Nursing_Education_Platform_gpt\Nursing_Education_Platform_gpt\docs\FIRST REVIEW PPT Template.pptx'
prs = Presentation(pptx_path)

for slide_idx, slide in enumerate(prs.slides):
    slide_num = slide_idx + 1
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            rows = len(tbl.rows)
            cols = len(tbl.columns)
            print(f'\n=== Slide {slide_num}: Table rows={rows} cols={cols} ===')
            for r in range(rows):
                row_data = []
                for c in range(cols):
                    cell = tbl.cell(r, c)
                    text = cell.text_frame.text.strip() if cell.text_frame else ''
                    row_data.append(f'[{text[:50]}]')
                print(f'  Row {r}: {" | ".join(row_data)}')
