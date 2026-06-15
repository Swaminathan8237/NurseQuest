from pptx import Presentation
import json

ppt_path = r"d:\E0224033\Projeccts\Nursing_Education_Platform_gpt\Nursing_Education_Platform_gpt\docs\FIRST REVIEW PPT Template.pptx"
prs = Presentation(ppt_path)

result = []
result.append(f"Slide dimensions: {prs.slide_width} x {prs.slide_height}")
result.append(f"Total slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    result.append(f"\n=== Slide {i+1} (Layout: {slide.slide_layout.name}) ===")
    for j, shape in enumerate(slide.shapes):
        result.append(f"  Shape {j}: name='{shape.name}', type={shape.shape_type}")
        result.append(f"    left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}")
        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
            result.append(f"    placeholder_idx={shape.placeholder_format.idx}")
        if shape.has_text_frame:
            for k, para in enumerate(shape.text_frame.paragraphs):
                txt = para.text
                font_info = ""
                if para.runs:
                    r = para.runs[0]
                    font_info = f" [font={r.font.name}, size={r.font.size}, bold={r.font.bold}]"
                result.append(f"    P{k}: '{txt[:100]}'{font_info}")

result.append("\n=== Layouts ===")
for i, layout in enumerate(prs.slide_layouts):
    result.append(f"  Layout {i}: '{layout.name}'")
    for ph in layout.placeholders:
        result.append(f"    PH idx={ph.placeholder_format.idx} name='{ph.name}'")

out_path = r"d:\E0224033\Projeccts\Nursing_Education_Platform_gpt\Nursing_Education_Platform_gpt\ppt_structure.txt"
with open(out_path, "w", encoding="utf-8") as f:
    f.write("\n".join(result))

print("Done! Output written to ppt_structure.txt")
